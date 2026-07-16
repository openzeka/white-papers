from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from slides.base import BaseSlide, Colors

def build(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    BaseSlide.set_background(slide)
    BaseSlide.add_title(slide, "İlk Karar: Eğitim mi? Inference mı?")
    BaseSlide.add_subtitle(slide, "Donanım ihtiyacını belirleyen en önemli soru")

    cx = Inches(6.666)

    d1 = BaseSlide.add_shape(slide, MSO_SHAPE.DIAMOND, cx - Inches(1.8), Inches(1.8), Inches(3.6), Inches(1.4),
                             fill_color=Colors.ACCENT_ORANGE, text="Eğitim yapılacak mı?", font_size=Pt(13), font_color=Colors.BG_DARK)

    lbl_yes = BaseSlide.add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(3.8), Inches(2.0), Inches(0.6),
                                   fill_color=Colors.BG_CARD, text="EVET", font_size=Pt(14), font_color=Colors.ACCENT_BLUE)
    lbl_yes.line.color.rgb = Colors.ACCENT_BLUE
    lbl_yes.line.width = Pt(1.5)

    lbl_no = BaseSlide.add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(3.8), Inches(2.0), Inches(0.6),
                                 fill_color=Colors.BG_CARD, text="HAYIR", font_size=Pt(14), font_color=Colors.ACCENT_BLUE)
    lbl_no.line.color.rgb = Colors.ACCENT_BLUE
    lbl_no.line.width = Pt(1.5)

    box_w = Inches(2.95)
    box_h = Inches(1.6)
    box_top = Inches(5.0)

    pre = BaseSlide.add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), box_top, box_w, box_h,
                              fill_color=Colors.BG_CARD, text="Pre-training (sıfırdan)\n8-32x B200/B300 (1-4 takım)\nÇok Yüksek VRAM", font_size=Pt(11))
    pre.line.color.rgb = Colors.ACCENT_ORANGE
    pre.line.width = Pt(1.5)

    fft = BaseSlide.add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), box_top, box_w, box_h,
                               fill_color=Colors.BG_CARD, text="Full Fine-tune\n1-8x H200 (ZeRO-3)\nYüksek VRAM", font_size=Pt(11))
    fft.line.color.rgb = Colors.ACCENT_ORANGE
    fft.line.width = Pt(1.5)

    lora = BaseSlide.add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), box_top, box_w, box_h,
                                fill_color=Colors.BG_CARD, text="Fine-tune (LoRA/QLoRA)\n1x L40S / RTX PRO 6000 / H200\nOrta/Düşük VRAM", font_size=Pt(11))
    lora.line.color.rgb = Colors.ACCENT_ORANGE
    lora.line.width = Pt(1.5)

    inf = BaseSlide.add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.0), box_top, Inches(3.0), box_h,
                               fill_color=Colors.BG_CARD, text="Inference\n1x H200 / 8x B200 (takım)\nKV cache → ek VRAM kritik", font_size=Pt(11))
    inf.line.color.rgb = Colors.ACCENT_ORANGE
    inf.line.width = Pt(1.5)

    BaseSlide.add_arrow_connector(slide, cx - Inches(1.0), Inches(3.2), Inches(5.0), Inches(3.8))
    BaseSlide.add_arrow_connector(slide, cx + Inches(1.0), Inches(3.2), Inches(11.5), Inches(3.8))

    BaseSlide.add_arrow_connector(slide, Inches(4.5), Inches(4.4), Inches(1.75), Inches(5.0))
    BaseSlide.add_arrow_connector(slide, Inches(5.0), Inches(4.4), Inches(4.95), Inches(5.0))
    BaseSlide.add_arrow_connector(slide, Inches(5.5), Inches(4.4), Inches(8.15), Inches(5.0))

    BaseSlide.add_arrow_connector(slide, Inches(11.5), Inches(4.4), Inches(11.5), Inches(5.0))

    BaseSlide.add_notes(slide, "İlk ve en önemli karar: Eğitim mi inference mı? B200/B300 8'li takım olarak satılır, bu bütçeyi ciddi etkiler. Pre-training (sıfırdan) en ağır senaryo — yalnızca kendi temel modelini eğiten kurumlar içindir. Full fine-tune tüm parametreleri günceller, H200 sınıfı çoklu-GPU ister. Fine-tune (LoRA/QLoRA) için tekli GPU'lar (L40S, RTX PRO 6000, H200) yeterli. Inference için düşük trafikte H200, yüksek trafikte B200 takım gerekli.")
