from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


class Colors:
    BG_DARK = RGBColor(0x0F, 0x17, 0x29)
    TEXT_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
    ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
    ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)
    BG_CARD = RGBColor(0x1E, 0x29, 0x3B)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT_DIM = RGBColor(0x94, 0xA3, 0xB8)


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


class BaseSlide:
    @staticmethod
    def set_background(slide, color=Colors.BG_DARK):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    @staticmethod
    def add_title(slide, text, left=Inches(0.8), top=Inches(0.4), width=Inches(11.7), height=Inches(0.9)):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = Colors.TEXT_LIGHT
        p.font.name = "Calibri"
        return txBox

    @staticmethod
    def add_subtitle(slide, text, left=Inches(0.8), top=Inches(1.3), width=Inches(11.7), height=Inches(0.6)):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = Colors.ACCENT_ORANGE
        p.font.name = "Calibri"
        return txBox

    @staticmethod
    def add_body(slide, text, left=Inches(0.8), top=Inches(2.0), width=Inches(11.7), height=Inches(4.5)):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text.split("\n")):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = Colors.TEXT_LIGHT
            p.font.name = "Calibri"
            p.space_after = Pt(6)
        return txBox

    @staticmethod
    def add_bullet_list(slide, items, left=Inches(0.8), top=Inches(2.2), width=Inches(11.7), height=Inches(4.5), font_size=Pt(18)):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = font_size
            p.font.color.rgb = Colors.TEXT_LIGHT
            p.font.name = "Calibri"
            p.space_after = Pt(8)
            p.level = 0
        return txBox

    @staticmethod
    def add_shape(slide, shape_type, left, top, width, height, fill_color=Colors.BG_CARD, text="", font_size=Pt(14), font_color=Colors.TEXT_LIGHT):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = font_size
            p.font.color.rgb = font_color
            p.font.name = "Calibri"
            p.font.bold = True
            tf.paragraphs[0].space_before = Pt(0)
            tf.paragraphs[0].space_after = Pt(0)
        return shape

    @staticmethod
    def add_connector(slide, start_x, start_y, end_x, end_y, color=Colors.ACCENT_ORANGE, width=Pt(2)):
        connector = slide.shapes.add_connector(
            1, start_x, start_y, end_x, end_y
        )
        connector.line.color.rgb = color
        connector.line.width = width
        return connector

    @staticmethod
    def add_arrow_connector(slide, start_x, start_y, end_x, end_y, color=Colors.ACCENT_ORANGE, width=Pt(2)):
        connector = slide.shapes.add_connector(
            1, start_x, start_y, end_x, end_y
        )
        connector.line.color.rgb = color
        connector.line.width = width
        from pptx.oxml.ns import qn
        ln = connector.line._ln
        tailEnd = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tailEnd)
        return connector

    @staticmethod
    def add_table(slide, rows, cols, data, left=Inches(0.5), top=Inches(2.2), width=Inches(12.3), height=Inches(4.5), header_color=Colors.ACCENT_ORANGE, font_size=Pt(12)):
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.text = str(data[r][c])
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = font_size
                    paragraph.font.name = "Calibri"
                    if r == 0:
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = Colors.BG_DARK
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = header_color
                    else:
                        paragraph.font.color.rgb = Colors.TEXT_LIGHT
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = Colors.BG_CARD
        return table_shape

    @staticmethod
    def add_notes(slide, text):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = text

    @staticmethod
    def add_card(slide, title, body_text, left, top, width=Inches(3.5), height=Inches(2.5)):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = Colors.BG_CARD
        shape.line.color.rgb = Colors.ACCENT_ORANGE
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = Colors.ACCENT_ORANGE
        p.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.text = body_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = Colors.TEXT_LIGHT
        p2.font.name = "Calibri"
        p2.space_before = Pt(8)
        return shape
